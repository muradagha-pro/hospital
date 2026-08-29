from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parents[1]
LOGO_PATH = BASE_DIR / "alhourani.jpeg"
OUTPUT_PATH = BASE_DIR / "Hospital_Management_Presentation_AR.pptx"


BG_COLOR = RGBColor(246, 244, 239)
TITLE_COLOR = RGBColor(30, 74, 70)
ACCENT_COLOR = RGBColor(46, 110, 103)
TEXT_COLOR = RGBColor(31, 46, 44)
MUTED_COLOR = RGBColor(108, 122, 120)
WHITE = RGBColor(255, 255, 255)


def style_text_run(run, size, bold=False, color=TEXT_COLOR):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_logo(slide, left=Inches(0.25), top=Inches(0.2), width=Inches(1.0)):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), left, top, width=width)


def add_footer(slide, text="نظام إدارة الاستدعاء والكافتيريا - عرض للإدارة"):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    style_text_run(run, 12, color=MUTED_COLOR)


def add_title(slide, title_text, subtitle_text=None):
    title_box = slide.shapes.add_textbox(Inches(1.4), Inches(0.7), Inches(11.3), Inches(1.2))
    tf = title_box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = title_text
    style_text_run(run, 34, bold=True, color=TITLE_COLOR)

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.RIGHT
        run2 = p2.add_run()
        run2.text = subtitle_text
        style_text_run(run2, 18, color=MUTED_COLOR)


def add_content_card(slide):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.15),
        Inches(1.75),
        Inches(11.0),
        Inches(4.85),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(228, 224, 214)
    card.line.width = Pt(1.3)
    return card


def add_bullets(slide, items, title=None):
    text_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(9.8), Inches(4.4))
    tf = text_box.text_frame
    tf.clear()

    if title:
        p_title = tf.paragraphs[0]
        p_title.alignment = PP_ALIGN.RIGHT
        run_title = p_title.add_run()
        run_title.text = title
        style_text_run(run_title, 22, bold=True, color=ACCENT_COLOR)
    else:
        p_title = tf.paragraphs[0]
        p_title.alignment = PP_ALIGN.RIGHT
        run_title = p_title.add_run()
        run_title.text = ""

    for item in items:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = f"• {item}"
        style_text_run(run, 20)


def add_side_visual(slide):
    visual = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.5),
        Inches(2.1),
        Inches(1.45),
        Inches(1.45),
    )
    visual.fill.solid()
    visual.fill.fore_color.rgb = RGBColor(234, 242, 233)
    visual.line.color.rgb = ACCENT_COLOR
    visual.line.width = Pt(1.2)

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(10.65), Inches(2.25), width=Inches(1.15))


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_logo(slide, left=Inches(0.35), top=Inches(0.25), width=Inches(1.25))

    add_title(
        slide,
        "عرض نظام مشفى الحوراني الرقمي",
        "منصة موحدة لاستدعاء التمريض وإدارة الكافتيريا والتحليلات الإدارية",
    )

    banner = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.3),
        Inches(2.5),
        Inches(10.2),
        Inches(1.3),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = ACCENT_COLOR
    banner.line.color.rgb = ACCENT_COLOR

    tf = banner.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "جاهز للتشغيل في بيئة المستشفى مع دعم التوسع المرحلي"
    style_text_run(run, 24, bold=True, color=WHITE)

    info = slide.shapes.add_textbox(Inches(1.3), Inches(4.2), Inches(10.4), Inches(2.0))
    info_tf = info.text_frame
    info_tf.clear()
    p1 = info_tf.paragraphs[0]
    p1.alignment = PP_ALIGN.RIGHT
    r1 = p1.add_run()
    r1.text = "الفئة المستهدفة: الإدارة العليا ومدراء التشغيل والجودة"
    style_text_run(r1, 20)

    p2 = info_tf.add_paragraph()
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = "الهدف: تسريع الاستجابة للمرضى ورفع كفاءة المتابعة واتخاذ القرار"
    style_text_run(r2, 20)

    add_footer(slide)


def add_content_slide(prs, main_title, section_title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_logo(slide)
    add_title(slide, main_title)
    add_content_card(slide)
    add_bullets(slide, bullets, title=section_title)
    add_side_visual(slide)
    add_footer(slide)


def build_presentation():
    if not LOGO_PATH.exists():
        raise FileNotFoundError(f"Logo image not found: {LOGO_PATH}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_content_slide(
        prs,
        "لماذا هذا النظام؟",
        "التحديات التي يعالجها",
        [
            "تقليل تأخر استجابة التمريض للحالات غير الطارئة.",
            "إلغاء الاعتماد على النداء الشفهي أو التنسيق اليدوي.",
            "تتبع لحظي لحالة كل طلب من الإرسال حتى الإغلاق.",
            "تجميع البيانات التشغيلية في لوحة واضحة للإدارة.",
        ],
    )

    add_content_slide(
        prs,
        "نظرة عامة على المكونات",
        "الواجهات الرئيسية",
        [
            "واجهة المريض: استدعاء ممرضة + متابعة الحالة + إرسال اقتراح/شكوى.",
            "واجهة الممرضة: استقبال الطلبات، الاستلام، الإنهاء، وتوثيق الملاحظات.",
            "لوحة الإدارة: مؤشرات الأداء، متوسطات الاستجابة، وملاحظات التمريض.",
            "نظام الكافتيريا: طلبات المرضى + إدارة المنتجات + تتبع التسليم.",
        ],
    )

    add_content_slide(
        prs,
        "رحلة المريض",
        "خطوات الاستخدام",
        [
            "المريض يمسح QR الغرفة ويفتح صفحة الخدمة مباشرة.",
            "زر استدعاء واضح يرسل الطلب فوريا للقسم المختص.",
            "المريض يرى حالة الطلب: مرسل / مستلم / مكتمل.",
            "يمكنه إرسال شكوى أو اقتراح للإدارة من نفس الرحلة.",
        ],
    )

    add_content_slide(
        prs,
        "رحلة التمريض",
        "آلية العمل اليومية",
        [
            "كل ممرضة تدخل من رابط القسم الخاص بها وتستقبل الحالات الفعلية فقط.",
            "تنبيهات صوتية واهتزازية لتقليل أي طلب فائت.",
            "إجراءات مباشرة: استلام الطلب ثم إنهاؤه بعد تقديم الخدمة.",
            "سجل يومي بالملاحظات يضمن التوثيق والمتابعة والتحسين.",
        ],
    )

    add_content_slide(
        prs,
        "لوحة الإدارة والتحليلات",
        "ما الذي يهم الإدارة؟",
        [
            "إجمالي الطلبات المنفذة والمعلقة لكل فترة زمنية.",
            "متوسط زمن الاستجابة على مستوى المستشفى ولكل قسم.",
            "عرض مركزي لملاحظات الممرضات وشكاوى المرضى.",
            "دعم قرارات سريعة لتحسين التوزيع والكوادر والجودة.",
        ],
    )

    add_content_slide(
        prs,
        "إدارة الأقسام الفعلية في النظام",
        "البيانات الحالية المعتمدة للعرض",
        [
            "النسائية - رمز القسم: dept-1787470659022 - الغرف: 300 إلى 399.",
            "قثطرة - رمز القسم: dept-1787485105394 - الغرف: 200 إلى 299.",
            "الخمسميات - رمز القسم: dept-1787486243729 - الغرف: 500 إلى 599.",
            "النظام يوجه كل غرفة تلقائيا إلى قسمها دون تدخل يدوي.",
        ],
    )

    add_content_slide(
        prs,
        "نظام الكافتيريا الذكي",
        "القيمة التشغيلية",
        [
            "المريض يطلب من نفس التجربة الرقمية دون اتصالات هاتفية.",
            "الكافتيريا تدير حالة الطلب: تم الاستلام / قيد التحضير / تم التسليم.",
            "متابعة شفافة للغرفة مع وقت تقديري وتفاصيل الفاتورة.",
            "تقليل الضغط على التمريض في الطلبات غير الطبية.",
        ],
    )

    add_content_slide(
        prs,
        "الأمن والجاهزية",
        "الوضع الحالي والخطوات التالية",
        [
            "النسخة الحالية مناسبة للتجربة التشغيلية الداخلية.",
            "قبل الإطلاق الرسمي: تفعيل تسجيل الدخول لصفحات الإدارة.",
            "تشديد قواعد قاعدة البيانات حسب الصلاحيات والأدوار.",
            "تفعيل سجل تدقيق إداري للتتبع والامتثال.",
        ],
    )

    add_content_slide(
        prs,
        "خطة التنفيذ المقترحة",
        "مراحل اعتماد الإدارة",
        [
            "مرحلة 1: تشغيل تجريبي في الأقسام الثلاثة الحالية لمدة أسبوعين.",
            "مرحلة 2: قياس زمن الاستجابة ورضا المرضى ومعدلات الإغلاق.",
            "مرحلة 3: تعميم تدريجي على بقية الأقسام مع تدريب المستخدمين.",
            "مرحلة 4: اعتماد نهائي وربط مؤشرات الجودة بالتقارير الشهرية.",
        ],
    )

    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


if __name__ == "__main__":
    result = build_presentation()
    print(f"Presentation generated: {result}")

